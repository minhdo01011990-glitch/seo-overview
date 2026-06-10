"""Tạo báo cáo SEO: HTML, DOCX, PPTX, PNG và upload lên Google Drive."""

from __future__ import annotations

import base64
import io
import json
import mimetypes
import os
import re
from pathlib import Path
from typing import Optional

# Templates dir: src/core/ → ../../templates
_TEMPLATES_DIR = Path(__file__).parent.parent.parent / "templates"


# ── HTML ──────────────────────────────────────────────────────────────────────

def generate_html(
    context: dict,
    output_path: str,
    templates_dir: Optional[str] = None,
) -> str:
    """Render report.html.j2 với Jinja2 và lưu ra file."""
    from jinja2 import Environment, FileSystemLoader

    tdir = templates_dir or str(_TEMPLATES_DIR)
    env = Environment(loader=FileSystemLoader(tdir), autoescape=True)
    template = env.get_template("report.html.j2")
    html = template.render(report=context)

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(html, encoding="utf-8")
    return str(out)


# ── DOCX ──────────────────────────────────────────────────────────────────────

def _add_section_to_docx(doc, section: dict) -> None:
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Inches

    doc.add_heading(section["title"], level=1)

    for insight in section.get("key_insights", []):
        doc.add_paragraph(insight, style="List Bullet")

    html = section.get("content_html", "")

    # Split by <h3> and <table> blocks; handle other content as plain text
    for part in re.split(
        r'(<h3[^>]*>.*?</h3>|<table[^>]*>.*?</table>)',
        html,
        flags=re.DOTALL,
    ):
        if re.match(r'<h3', part):
            text = re.sub(r'<[^>]+>', '', part).strip()
            if text:
                doc.add_heading(text, level=2)

        elif re.match(r'<table', part):
            rows: list[list[str]] = []
            for tr_m in re.finditer(r'<tr[^>]*>(.*?)</tr>', part, re.DOTALL):
                cells = re.findall(r'<t[dh][^>]*>(.*?)</t[dh]>', tr_m.group(1), re.DOTALL)
                row = [re.sub(r'<[^>]+>', '', c).strip() for c in cells]
                if row:
                    rows.append(row)
            if rows:
                n_cols = max(len(r) for r in rows)
                try:
                    tbl = doc.add_table(rows=len(rows), cols=n_cols, style="Table Grid")
                except KeyError:
                    tbl = doc.add_table(rows=len(rows), cols=n_cols)
                for i, row in enumerate(rows):
                    for j, cell_text in enumerate(row[:n_cols]):
                        tbl.cell(i, j).text = cell_text
                doc.add_paragraph()

        else:
            # Skip kpi-grid markup (already summarised in key_insights bullets)
            if "kpi-grid" not in part and "kpi-card" not in part:
                text = re.sub(r'<[^>]+>', ' ', part)
                text = re.sub(r'\s+', ' ', text).strip()
                if text:
                    doc.add_paragraph(text)

    # Embed charts as inline images
    for chart in section.get("charts", []):
        try:
            img_data = base64.b64decode(chart["image_b64"])
            doc.add_picture(io.BytesIO(img_data), width=Inches(6))
            if chart.get("title"):
                p = doc.add_paragraph(chart["title"])
                p.alignment = WD_ALIGN_PARAGRAPH.CENTER
        except Exception as chart_exc:
            doc.add_paragraph(f"[Chart '{chart.get('title', '?')}' không thể render: {chart_exc}]")

    doc.add_paragraph()


def generate_docx(context: dict, output_path: str) -> str:
    """Tạo file DOCX từ report context."""
    from docx import Document

    doc = Document()

    doc.add_heading(context["title"], level=0)
    doc.add_paragraph(f"Domain: {context['my_domain']}")
    doc.add_paragraph(f"Ngày tạo: {context['generated_at']}")
    doc.add_page_break()

    for section in context["sections"]:
        _add_section_to_docx(doc, section)

    if context.get("missing_sections"):
        doc.add_heading("Sections không có dữ liệu", level=1)
        for ms in context["missing_sections"]:
            doc.add_paragraph(f"{ms['title']}: {ms['reason']}", style="List Bullet")

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    doc.save(str(out))
    return str(out)


# ── PPTX ──────────────────────────────────────────────────────────────────────

def _blank_layout(prs):
    """Trả về Blank slide layout (layout 6 trong default theme, fallback về cuối)."""
    try:
        return prs.slide_layouts[6]
    except IndexError:
        return prs.slide_layouts[-1]


def _add_pptx_title_textbox(slide, text: str, x, y, w, h, font_size: int = 22) -> None:
    from pptx.util import Pt

    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True


def _add_pptx_section_slide(prs, section: dict) -> None:
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(_blank_layout(prs))
    W = prs.slide_width
    H = prs.slide_height

    # Title bar — full width, top
    _add_pptx_title_textbox(
        slide, section["title"],
        Inches(0.3), Inches(0.15),
        W - Inches(0.6), Inches(0.75),
        font_size=22,
    )

    insights = section.get("key_insights", [])
    charts = section.get("charts", [])
    content_top = Inches(1.05)
    content_h = H - Inches(1.3)

    if charts:
        # Left 38%: key insights as bullets
        txt_w = int(W * 0.38)
        tb = slide.shapes.add_textbox(Inches(0.3), content_top, txt_w, content_h)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, insight in enumerate(insights[:10]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {insight}"
            p.font.size = Pt(10)

        # Right 57%: first chart
        chart_left = int(W * 0.41)
        chart_w = W - chart_left - Inches(0.2)
        img_data = base64.b64decode(charts[0]["image_b64"])
        slide.shapes.add_picture(
            io.BytesIO(img_data), chart_left, content_top, chart_w, content_h
        )
    else:
        # Full width: bullets
        tb = slide.shapes.add_textbox(Inches(0.3), content_top, W - Inches(0.6), content_h)
        tf = tb.text_frame
        tf.word_wrap = True
        for i, insight in enumerate(insights):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {insight}"
            p.font.size = Pt(12)


def generate_pptx(context: dict, output_path: str) -> str:
    """Tạo file PPTX từ report context."""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()

    # Title slide (layout 0)
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    if slide.shapes.title:
        slide.shapes.title.text = context["title"]
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            ph.text = (
                f"Domain: {context['my_domain']}\n{context['generated_at']}"
            )
            break

    # One slide per section
    for section in context["sections"]:
        _add_pptx_section_slide(prs, section)

    # Missing sections slide
    if context.get("missing_sections"):
        slide = prs.slides.add_slide(_blank_layout(prs))
        W = prs.slide_width
        H = prs.slide_height
        _add_pptx_title_textbox(
            slide, "Sections không có dữ liệu",
            Inches(0.3), Inches(0.15),
            W - Inches(0.6), Inches(0.75),
            font_size=18,
        )
        tb = slide.shapes.add_textbox(Inches(0.3), Inches(1.05), W - Inches(0.6), H - Inches(1.3))
        tf = tb.text_frame
        tf.word_wrap = True
        for i, ms in enumerate(context["missing_sections"]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = f"• {ms['title']}: {ms['reason']}"
            p.font.size = Pt(12)

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return str(out)


# ── PNG ───────────────────────────────────────────────────────────────────────

def generate_png(context: dict, output_dir: str) -> list[str]:
    """
    Lưu charts mỗi section thành 1 file PNG (composite dọc nếu nhiều chart).
    Chỉ tạo file cho sections có ít nhất 1 chart.
    Trả về list absolute paths.
    """
    from PIL import Image

    out_dir = Path(output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)

    paths: list[str] = []
    for section in context["sections"]:
        charts = section.get("charts", [])
        if not charts:
            continue

        images = []
        for c in charts:
            try:
                images.append(Image.open(io.BytesIO(base64.b64decode(c["image_b64"]))))
            except Exception:
                continue
        if not images:
            continue

        if len(images) == 1:
            composite = images[0].convert("RGB")
        else:
            total_w = max(img.width for img in images)
            total_h = sum(img.height for img in images)
            composite = Image.new("RGB", (total_w, total_h), (255, 255, 255))
            y_off = 0
            for img in images:
                rgb = img.convert("RGB")
                composite.paste(rgb, ((total_w - rgb.width) // 2, y_off))
                y_off += img.height

        out_path = out_dir / f"{section['id']}.png"
        composite.save(str(out_path), format="PNG")
        paths.append(str(out_path))

    return paths


# ── Google Drive upload ───────────────────────────────────────────────────────

def upload_to_drive(file_paths: list[str], drive_folder_url: str) -> list[dict]:
    """Upload files lên Google Drive folder và trả về danh sách {file, drive_id, url}."""
    folder_id_m = re.search(r'/folders/([a-zA-Z0-9_-]+)', drive_folder_url)
    if not folder_id_m:
        raise ValueError(
            f"Không thể extract folder ID từ URL: {drive_folder_url}\n"
            "URL phải dạng https://drive.google.com/drive/folders/<id>"
        )
    folder_id = folder_id_m.group(1)

    sa_json = os.environ.get("GOOGLE_SERVICE_ACCOUNT_JSON", "")
    if not sa_json:
        raise EnvironmentError(
            "GOOGLE_SERVICE_ACCOUNT_JSON chưa được set trong environment"
        )

    from google.oauth2 import service_account
    from googleapiclient.discovery import build
    from googleapiclient.http import MediaFileUpload

    if sa_json.strip().startswith("{"):
        creds_info = json.loads(sa_json)
    else:
        with open(sa_json, encoding="utf-8") as f:
            creds_info = json.load(f)

    creds = service_account.Credentials.from_service_account_info(
        creds_info,
        scopes=["https://www.googleapis.com/auth/drive.file"],
    )
    service = build("drive", "v3", credentials=creds, cache_discovery=False)

    results: list[dict] = []
    for file_path in file_paths:
        p = Path(file_path)
        if not p.exists():
            results.append({"file": p.name, "error": "File không tồn tại"})
            continue

        mime_type, _ = mimetypes.guess_type(str(p))
        mime_type = mime_type or "application/octet-stream"

        file_obj = service.files().create(
            body={"name": p.name, "parents": [folder_id]},
            media_body=MediaFileUpload(str(p), mimetype=mime_type, resumable=False),
            fields="id,name,webViewLink",
        ).execute()

        results.append({
            "file": p.name,
            "drive_id": file_obj.get("id"),
            "url": file_obj.get("webViewLink"),
        })

    return results
